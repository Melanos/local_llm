"""
Document Image Extractor - Extract embedded images from documents
"""

import os
import sys
from pathlib import Path
from typing import List, Dict, Tuple, Optional
import zipfile
import fitz  # PyMuPDF for PDF
from pptx import Presentation
from docx import Document
import shutil
from PIL import Image
import io

# Add project root to path
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from config import DEFAULT_CONFIG, DOCUMENTS_DIR


class DocumentImageExtractor:
    """Extract images from various document types"""
    
    def __init__(self, config: Dict = None):
        self.config = config or DEFAULT_CONFIG
        self.extracted_images_dir = Path("data/images/extracted")
        self.extracted_images_dir.mkdir(exist_ok=True)
    
    def extract_images_from_docx(self, docx_path: Path) -> List[Dict]:
        """Extract images from a DOCX file"""
        extracted_images = []
        
        try:
            # DOCX files are ZIP archives
            with zipfile.ZipFile(docx_path, 'r') as zip_file:
                # Look for images in word/media/ folder
                image_files = [f for f in zip_file.namelist() if f.startswith('word/media/')]
                
                for i, image_file in enumerate(image_files):
                    try:
                        # Extract image data
                        image_data = zip_file.read(image_file)
                        
                        # Determine file extension
                        original_ext = Path(image_file).suffix.lower()
                        if not original_ext:
                            # Try to determine from image data
                            try:
                                img = Image.open(io.BytesIO(image_data))
                                original_ext = f".{img.format.lower()}"
                            except:
                                original_ext = ".png"  # Default
                        
                        # Create filename
                        base_name = docx_path.stem
                        new_filename = f"{base_name}_img{i+1}{original_ext}"
                        output_path = self.extracted_images_dir / new_filename
                        
                        # Save image
                        with open(output_path, 'wb') as img_file:
                            img_file.write(image_data)
                        
                        extracted_images.append({
                            'source_document': docx_path.name,
                            'extracted_path': output_path,
                            'image_index': i + 1,
                            'original_name': Path(image_file).name
                        })
                        
                        print(f"📷 Extracted image {i+1} from {docx_path.name} -> {new_filename}")
                        
                    except Exception as e:
                        print(f"⚠️ Error extracting image {image_file}: {e}")
                        
        except Exception as e:
            print(f"❌ Error processing DOCX {docx_path.name}: {e}")
        
        return extracted_images
    
    def extract_images_from_pdf(self, pdf_path: Path) -> List[Dict]:
        """Extract images from a PDF file"""
        extracted_images = []
        
        try:
            pdf_document = fitz.open(pdf_path)
            image_count = 0
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_document, xref)
                        
                        # Skip if not RGB or GRAY
                        if pix.n - pix.alpha < 4:
                            image_count += 1
                            
                            # Determine format
                            if pix.n - pix.alpha == 1:  # Grayscale
                                img_ext = ".png"
                            else:  # RGB
                                img_ext = ".png"
                            
                            # Create filename
                            base_name = pdf_path.stem
                            new_filename = f"{base_name}_p{page_num+1}_img{img_index+1}{img_ext}"
                            output_path = self.extracted_images_dir / new_filename
                            
                            # Save image
                            if pix.n - pix.alpha == 1:
                                pix.save(str(output_path))
                            else:
                                pix.save(str(output_path))
                            
                            extracted_images.append({
                                'source_document': pdf_path.name,
                                'extracted_path': output_path,
                                'page_number': page_num + 1,
                                'image_index': img_index + 1,
                                'total_image_number': image_count
                            })
                            
                            print(f"📷 Extracted image from {pdf_path.name} page {page_num+1} -> {new_filename}")
                        
                        pix = None  # Free memory
                        
                    except Exception as e:
                        print(f"⚠️ Error extracting image from page {page_num+1}: {e}")
            
            pdf_document.close()
            
        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path.name}: {e}")
        
        return extracted_images
    
    def extract_images_from_pptx(self, pptx_path: Path) -> List[Dict]:
        """Extract images from a PowerPoint file"""
        extracted_images = []
        
        try:
            prs = Presentation(pptx_path)
            image_count = 0
            
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_count += 1
                            image = shape.image
                            
                            # Get image data
                            image_bytes = image.blob
                            
                            # Determine extension from content type or filename
                            ext = ".png"  # Default
                            if hasattr(image, 'content_type'):
                                if 'jpeg' in image.content_type.lower():
                                    ext = ".jpg"
                                elif 'png' in image.content_type.lower():
                                    ext = ".png"
                            
                            # Create filename
                            base_name = pptx_path.stem
                            new_filename = f"{base_name}_slide{slide_num+1}_img{image_count}{ext}"
                            output_path = self.extracted_images_dir / new_filename
                            
                            # Save image
                            with open(output_path, 'wb') as img_file:
                                img_file.write(image_bytes)
                            
                            extracted_images.append({
                                'source_document': pptx_path.name,
                                'extracted_path': output_path,
                                'slide_number': slide_num + 1,
                                'image_index': image_count
                            })
                            
                            print(f"📷 Extracted image from {pptx_path.name} slide {slide_num+1} -> {new_filename}")
                            
                        except Exception as e:
                            print(f"⚠️ Error extracting image from slide {slide_num+1}: {e}")
            
        except Exception as e:
            print(f"❌ Error processing PowerPoint {pptx_path.name}: {e}")
        
        return extracted_images
    
    def extract_images_from_document(self, doc_path: Path) -> List[Dict]:
        """Extract images from a single document"""
        suffix = doc_path.suffix.lower()
        
        if suffix == '.docx':
            return self.extract_images_from_docx(doc_path)
        elif suffix == '.pdf':
            return self.extract_images_from_pdf(doc_path)
        elif suffix == '.pptx':
            return self.extract_images_from_pptx(doc_path)
        else:
            print(f"⚠️ Image extraction not supported for {suffix} files")
            return []
    
    def extract_all_document_images(self, documents_dir: Path = None) -> Dict:
        """Extract images from all documents in the directory"""
        if documents_dir is None:
            documents_dir = DOCUMENTS_DIR
        
        documents_dir = Path(documents_dir)
        
        if not documents_dir.exists():
            print(f"❌ Documents directory not found: {documents_dir}")
            return {"processed": 0, "extracted": 0, "errors": 0}
        
        print(f"🔍 Looking for documents with embedded images in: {documents_dir}")
        
        # Find supported documents
        supported_extensions = ['*.docx', '*.pdf', '*.pptx']
        doc_files = []
        
        for pattern in supported_extensions:
            doc_files.extend(documents_dir.glob(pattern))
        
        if not doc_files:
            print(f"📁 No supported documents found")
            return {"processed": 0, "extracted": 0, "errors": 0}
        
        print(f"📚 Found {len(doc_files)} document(s) to check for images")
        
        total_extracted = 0
        processed = 0
        errors = 0
        all_extracted_images = []
        
        for doc_file in doc_files:
            try:
                print(f"\n📄 Checking: {doc_file.name}")
                extracted = self.extract_images_from_document(doc_file)
                
                if extracted:
                    print(f"✅ Extracted {len(extracted)} image(s) from {doc_file.name}")
                    total_extracted += len(extracted)
                    all_extracted_images.extend(extracted)
                else:
                    print(f"📷 No images found in {doc_file.name}")
                
                processed += 1
                
            except Exception as e:
                print(f"❌ Error processing {doc_file.name}: {e}")
                errors += 1
        
        # Summary
        print(f"\n✅ Image extraction complete!")
        print(f"📊 Documents processed: {processed}")
        print(f"🖼️ Total images extracted: {total_extracted}")
        print(f"❌ Errors: {errors}")
        
        if total_extracted > 0:
            print(f"📁 Extracted images saved to: {self.extracted_images_dir}")
            print(f"💡 Next steps:")
            print(f"   1. Review extracted images in {self.extracted_images_dir}")
            print(f"   2. Move relevant images to data/images/ for training")
            print(f"   3. Run: python train_images.py")
        
        return {
            "processed": processed,
            "extracted": total_extracted,
            "errors": errors,
            "extracted_images": all_extracted_images
        }


def main():
    """Main function for standalone execution"""
    print("🖼️ Document Image Extractor")
    print("=" * 50)
    
    extractor = DocumentImageExtractor()
    results = extractor.extract_all_document_images()
    
    if results["extracted"] > 0:
        print(f"\n🎯 Successfully extracted {results['extracted']} images!")
        print(f"📂 Check the extracted images in: data/images/extracted/")
    else:
        print("\n📷 No images found in documents")


if __name__ == "__main__":
    main()
