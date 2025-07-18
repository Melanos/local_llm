"""
Document Training Module - Process and add documents to RAG database
"""

import sys
import os
from pathlib import Path
from typing import List, Dict, Any
from docx import Document
import PyPDF2
import pdfplumber
from pptx import Presentation
from openpyxl import load_workbook
import uuid
from datetime import datetime

# Add project root to path
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.core.rag_engine import RAGEngine
from src.utils.document_image_extractor import DocumentImageExtractor
from config import DEFAULT_CONFIG, DOCUMENTS_DIR


class DocumentTrainer:
    """Handle document processing and training"""
    
    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or DEFAULT_CONFIG
        self.rag_engine = RAGEngine(config)
        self.image_extractor = DocumentImageExtractor(config)
        self.auto_extract_images = self.config.get("features", {}).get("auto_extract_images", False)
    
    def extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text content from a Word document"""
        try:
            doc = Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content)
            
        except Exception as e:
            print(f"❌ Error extracting text from {file_path}: {e}")
            return None
    
    def extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text content from a PDF file"""
        try:
            # Try pdfplumber first (better text extraction)
            try:
                with pdfplumber.open(file_path) as pdf:
                    text_content = []
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_content.append(page_text.strip())
                    
                    if text_content:
                        return "\n".join(text_content)
            except Exception:
                pass
            
            # Fallback to PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = []
                
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text.strip())
                
                return "\n".join(text_content)
                
        except Exception as e:
            print(f"❌ Error extracting text from PDF {file_path}: {e}")
            return None
    
    def extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text content from a plain text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-8-sig', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read().strip()
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, read as binary and decode with errors='ignore'
            with open(file_path, 'rb') as file:
                return file.read().decode('utf-8', errors='ignore').strip()
                
        except Exception as e:
            print(f"❌ Error extracting text from {file_path}: {e}")
            return None
    
    def extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text content from a PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.append(" ".join(slide_text))
            
            return "\n".join(text_content)
            
        except Exception as e:
            print(f"❌ Error extracting text from PowerPoint {file_path}: {e}")
            return None
    
    def extract_text_from_xlsx(self, file_path: Path) -> str:
        """Extract text content from an Excel file"""
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"Sheet: {sheet_name}"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value).strip())
                    
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                if len(sheet_text) > 1:  # More than just the sheet name
                    text_content.append("\n".join(sheet_text))
            
            workbook.close()
            return "\n\n".join(text_content)
            
        except Exception as e:
            print(f"❌ Error extracting text from Excel {file_path}: {e}")
            return None
    
    def get_file_extractor(self, file_path: Path):
        """Get the appropriate text extractor for a file based on its extension"""
        suffix = file_path.suffix.lower()
        
        extractors = {
            '.docx': self.extract_text_from_docx,
            '.pdf': self.extract_text_from_pdf,
            '.txt': self.extract_text_from_txt,
            '.md': self.extract_text_from_txt,
            '.py': self.extract_text_from_txt,
            '.js': self.extract_text_from_txt,
            '.html': self.extract_text_from_txt,
            '.css': self.extract_text_from_txt,
            '.json': self.extract_text_from_txt,
            '.xml': self.extract_text_from_txt,
            '.csv': self.extract_text_from_txt,
            '.pptx': self.extract_text_from_pptx,
            '.xlsx': self.extract_text_from_xlsx,
        }
        
        return extractors.get(suffix)
    
    def is_supported_file(self, file_path: Path) -> bool:
        """Check if the file type is supported"""
        return self.get_file_extractor(file_path) is not None
    
    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """Split text into overlapping chunks"""
        if chunk_size is None:
            chunk_size = self.config["chunking"]["chunk_size"]
        if overlap is None:
            overlap = self.config["chunking"]["overlap"]
            
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to end at a sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                boundary = max(last_period, last_newline)
                
                if boundary > start + chunk_size // 2:  # Only if boundary is reasonable
                    end = start + boundary + 1
                    chunk = text[start:end]
            
            chunks.append(chunk.strip())
            start = end - overlap
            
            if end >= len(text):
                break
                
        return chunks
    
    def process_document(self, file_path: Path) -> bool:
        """Process a single document and add to database"""
        print(f"📄 Processing: {file_path.name}")
        
        if not file_path.exists():
            print(f"❌ File not found: {file_path}")
            return False

        # Extract embedded images if enabled and supported
        if self.auto_extract_images and file_path.suffix.lower() in ['.docx', '.pdf', '.pptx']:
            try:
                print(f"🖼️ Checking for embedded images in {file_path.name}...")
                extracted_images = self.image_extractor.extract_images_from_document(file_path)
                if extracted_images:
                    print(f"📷 Extracted {len(extracted_images)} image(s) from {file_path.name}")
                else:
                    print(f"📷 No images found in {file_path.name}")
            except Exception as e:
                print(f"⚠️ Error extracting images from {file_path.name}: {e}")

        # Extract text based on file type
        extractor = self.get_file_extractor(file_path)
        if not extractor:
            print(f"❌ Unsupported file type: {file_path.suffix}")
            return False
            
        text = extractor(file_path)
        
        if not text:
            print(f"❌ No text extracted from {file_path.name}")
            return False
        
        # Check for existing document to avoid duplicates
        existing_docs = self.rag_engine.search_documents(
            f"source:{file_path.name}", n_results=1
        )
        
        if existing_docs['documents'][0]:
            print(f"⚠️  Document {file_path.name} already exists in database")
            return False
        
        # Chunk the text
        chunks = self.chunk_text(text)
        print(f"📝 Created {len(chunks)} chunks")
        
        # Add chunks to database
        added_count = 0
        for i, chunk in enumerate(chunks):
            if len(chunk.strip()) < 50:  # Skip very small chunks
                continue
                
            metadata = {
                'source': file_path.name,
                'source_type': 'document',
                'file_type': file_path.suffix.lower(),
                'chunk_id': i,
                'chunk_count': len(chunks),
                'file_size': file_path.stat().st_size,
                'processed_date': datetime.now().isoformat()
            }
            
            doc_id = self.rag_engine.add_document(chunk, metadata)
            if doc_id:
                added_count += 1
        
        print(f"✅ Added {added_count} chunks from {file_path.name}")
        return True
    
    def train_documents(self, directory: Path = None) -> Dict[str, Any]:
        """Train on all documents in the specified directory"""
        if directory is None:
            directory = DOCUMENTS_DIR
        
        directory = Path(directory)
        
        if not directory.exists():
            print(f"❌ Directory not found: {directory}")
            return {"processed": 0, "errors": 0}
        
        print(f"🔍 Looking for documents in: {directory}")
        
        # Find all supported document files
        doc_files = []
        supported_extensions = [
            '*.docx', '*.pdf', '*.txt', '*.md', '*.py', '*.js', 
            '*.html', '*.css', '*.json', '*.xml', '*.csv', 
            '*.pptx', '*.xlsx'
        ]
        
        for pattern in supported_extensions:
            doc_files.extend(directory.glob(pattern))
        
        # Filter to only supported files
        doc_files = [f for f in doc_files if self.is_supported_file(f)]
        
        if not doc_files:
            print(f"📁 No supported document files found in {directory}")
            print("📝 Supported formats: .docx, .pdf, .txt, .md, .py, .js, .html, .css, .json, .xml, .csv, .pptx, .xlsx")
            return {"processed": 0, "errors": 0}
        
        print(f"📚 Found {len(doc_files)} document(s) to process")
        
        # Process each document
        processed = 0
        errors = 0
        
        for doc_file in doc_files:
            try:
                if self.process_document(doc_file):
                    processed += 1
                else:
                    errors += 1
            except Exception as e:
                print(f"❌ Error processing {doc_file.name}: {e}")
                errors += 1
        
        # Show results
        print(f"\n✅ Document training complete!")
        print(f"📊 Processed: {processed} documents")
        print(f"❌ Errors: {errors} documents")
        
        # Show database stats
        stats = self.rag_engine.get_stats()
        print(f"💾 Total documents in database: {stats['total_documents']}")
        
        return {
            "processed": processed,
            "errors": errors,
            "total_in_db": stats['total_documents']
        }


def main():
    """Main entry point for document training"""
    print("📚 Document Training Module")
    print("=" * 40)
    
    trainer = DocumentTrainer()
    
    # Check if documents directory exists and has files
    if not DOCUMENTS_DIR.exists():
        print(f"📁 Creating documents directory: {DOCUMENTS_DIR}")
        DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
        print(f"📝 Please add .docx files to: {DOCUMENTS_DIR}")
        return
    
    results = trainer.train_documents()
    
    if results["processed"] > 0:
        print(f"\n🎉 Successfully trained on {results['processed']} documents!")
        print("💬 You can now chat with your documents using: python chat.py")
    else:
        print(f"\n📝 No documents processed. Add .docx files to: {DOCUMENTS_DIR}")


if __name__ == "__main__":
    main()
